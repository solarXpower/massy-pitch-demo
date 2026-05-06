"""BiziBid Partner Presentation Generator v2 — 18 slides with phone mockups
Run: python generate_presentation.py
Outputs: BiziBid-Partner-Presentation.pptx + BiziBid-Partner-Presentation-PRINT.html
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors ───────────────────────────────────────────────────────────────────
GOLD       = RGBColor(0xC9, 0xA8, 0x4C)
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xF8, 0xF8, 0xF8)
MGREY      = RGBColor(0x4B, 0x55, 0x63)
LGREY      = RGBColor(0x9C, 0xA3, 0xAF)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
NAVY       = RGBColor(0x0B, 0x1F, 0x3A)
DKGOLD     = RGBColor(0x92, 0x70, 0x20)
BUZZ_DARK  = RGBColor(0x06, 0x12, 0x2A)
BUZZ_BLUE  = RGBColor(0x1E, 0x50, 0xD8)
BUZZ_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PHONE_BODY = RGBColor(0x10, 0x10, 0x10)
PHONE_SCRN = RGBColor(0x0C, 0x0C, 0x1C)
APP_DARK   = RGBColor(0x0E, 0x2A, 0x44)
APP_CELL   = RGBColor(0x12, 0x12, 0x26)
NOTIF_GRN  = RGBColor(0x05, 0x96, 0x69)
NOTIF_RED  = RGBColor(0xC0, 0x20, 0x20)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
TOTAL = 18

# ── Core Helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        if line_w: sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def slide_bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def gold_bar(slide, t=0.92, h=0.07):
    add_rect(slide, 0, t, 13.33, h, fill=GOLD)

def page_num(slide, n):
    add_text(slide, f"{n} / {TOTAL}", 12.3, 7.1, 1.0, 0.35,
             size=9, color=LGREY, align=PP_ALIGN.RIGHT)

def logo_text(slide, l=0.4, t=0.18):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(2.5), Inches(0.55))
    tf = tb.text_frame; p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Bizi"
    r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = BLACK
    r2 = p.add_run(); r2.text = "Bid"
    r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = GOLD

def section_label(slide, text, l=0.5, t=1.05):
    add_text(slide, text.upper(), l, t, 8, 0.32, size=9, bold=True, color=GOLD)

def heading(slide, text, l=0.5, t=1.45, w=12.3, size=30, color=BLACK):
    add_text(slide, text, l, t, w, 0.85, size=size, bold=True, color=color)

def sub(slide, text, l=0.5, t=2.1, w=12.3, size=15, color=MGREY):
    add_text(slide, text, l, t, w, 0.45, size=size, color=color)

def feature_card(slide, icon, title, body, l, t, w=3.9, h=1.7,
                 bg=LIGHT, icon_color=GOLD, title_color=BLACK, body_color=MGREY):
    add_rect(slide, l, t, w, h, fill=bg)
    add_rect(slide, l, t, 0.05, h, fill=icon_color)
    add_text(slide, icon, l+0.12, t+0.12, 0.5, 0.45, size=20, color=icon_color)
    add_text(slide, title, l+0.65, t+0.12, w-0.8, 0.38, size=12, bold=True, color=title_color)
    add_text(slide, body,  l+0.12, t+0.58, w-0.25, h-0.7, size=10, color=body_color)

def stat_card(slide, number, label, l, t, w=2.2, h=1.1):
    add_rect(slide, l, t, w, h, fill=LIGHT)
    add_rect(slide, l, t, w, 0.04, fill=GOLD)
    add_text(slide, number, l, t+0.1, w, 0.5, size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, label,  l, t+0.6, w, 0.45, size=10, color=MGREY, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h, text_color=BLACK, text_size=13):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run(); r.text = item
        r.font.size = Pt(text_size); r.font.color.rgb = text_color
    return box

# ── Phone Mockup Helpers ──────────────────────────────────────────────────────

def draw_phone(slide, l, t, w=2.2):
    """Draw a phone frame. Returns (screen_l, screen_t, screen_w, screen_h)."""
    h = w * 2.08
    bevel  = w * 0.052
    top_bz = w * 0.092
    bot_bz = w * 0.082
    add_rect(slide, l, t, w, h, fill=PHONE_BODY)
    sl = l + bevel
    st = t + top_bz
    sw = w - bevel * 2
    sh = h - top_bz - bot_bz
    add_rect(slide, sl, st, sw, sh, fill=PHONE_SCRN)
    # speaker bar
    add_rect(slide, l + (w - sw*0.28)/2, t + top_bz*0.28, sw*0.28, top_bz*0.22, fill=LGREY)
    # home indicator bar
    add_rect(slide, l + (w - sw*0.24)/2, t + h - bot_bz*0.34, sw*0.24, bot_bz*0.14, fill=LGREY)
    return sl, st, sw, sh

def ph_fill(slide, sl, st, sw, sh, color=APP_DARK):
    add_rect(slide, sl, st, sw, sh, fill=color)

def ph_status(slide, sl, st, sw):
    """Minimal status bar. Returns y after."""
    bh = 0.105
    add_rect(slide, sl, st, sw, bh, fill=APP_DARK)
    add_text(slide, "9:41", sl+0.04, st+0.005, sw*0.3, bh,
             size=4.5, bold=True, color=WHITE)
    add_text(slide, "|||  100%", sl+sw*0.55, st+0.005, sw*0.43, bh,
             size=4, color=WHITE, align=PP_ALIGN.RIGHT)
    return st + bh

def ph_header(slide, sl, y, sw, text, bg=GOLD, fg=BLACK):
    """App top bar. Returns y after."""
    bh = 0.205
    add_rect(slide, sl, y, sw, bh, fill=bg)
    add_text(slide, text, sl+0.03, y+0.04, sw-0.06, bh-0.05,
             size=6.5, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return y + bh

def ph_label(slide, sl, y, sw, text, color=LGREY):
    """Small section label. Returns y after."""
    lh = 0.1
    add_text(slide, text, sl+0.04, y, sw-0.06, lh, size=3.8, color=color)
    return y + lh

def ph_field(slide, sl, y, sw, label, value="", fh=0.188):
    """Form input field. Returns y after."""
    add_rect(slide, sl, y, sw, fh, fill=APP_CELL)
    add_text(slide, label, sl+0.05, y+0.01, sw-0.08, 0.08, size=3.8, color=LGREY)
    add_text(slide, value, sl+0.05, y+0.09, sw-0.08, 0.09, size=5, bold=True, color=WHITE)
    add_rect(slide, sl, y+fh-0.007, sw, 0.007, fill=PHONE_BODY)
    return y + fh

def ph_notif(slide, sl, y, sw, title, body, badge=NOTIF_GRN, ts="now"):
    """Push notification row. Returns y after."""
    nh = 0.265
    add_rect(slide, sl, y, sw, nh, fill=APP_CELL)
    add_rect(slide, sl, y, 0.025, nh, fill=badge)
    add_text(slide, title, sl+0.05, y+0.02, sw-0.12, 0.10, size=4.5, bold=True, color=WHITE)
    add_text(slide, body,  sl+0.05, y+0.12, sw-0.12, 0.11, size=3.8, color=LGREY)
    add_text(slide, ts, sl+sw-0.14, y+0.02, 0.12, 0.08, size=3.5, color=LGREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, sl, y+nh, sw, 0.008, fill=PHONE_BODY)
    return y + nh + 0.008

def ph_btn(slide, sl, y, sw, text, bg=GOLD, fg=BLACK, bh=0.175):
    """Full-width button. Returns y after."""
    mx = sw * 0.09
    add_rect(slide, sl+mx, y, sw-mx*2, bh, fill=bg)
    add_text(slide, text, sl+mx, y+0.03, sw-mx*2, bh-0.04,
             size=5.5, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return y + bh + 0.025

def ph_bid_row(slide, sl, y, sw, user, amount, winning=False):
    """Bid history row. Returns y after."""
    rh = 0.148
    bg = RGBColor(0x18, 0x28, 0x40) if winning else APP_CELL
    add_rect(slide, sl, y, sw, rh, fill=bg)
    add_text(slide, user,   sl+0.04, y+0.03, sw*0.55, 0.09, size=3.8, color=WHITE if winning else LGREY)
    add_text(slide, amount, sl+sw*0.58, y+0.03, sw*0.4, 0.09, size=4.5, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    if winning:
        add_rect(slide, sl+sw-0.02, y, 0.02, rh, fill=GREEN)
    add_rect(slide, sl, y+rh, sw, 0.007, fill=PHONE_BODY)
    return y + rh + 0.007

def ph_store_key(slide, x, y, kw, kh, name, badge_color=None, badge_txt=""):
    """Store keypad button."""
    add_rect(slide, x, y, kw, kh, fill=APP_CELL)
    add_text(slide, name[0].upper(), x+kw*0.2, y+kh*0.1, kw*0.6, kh*0.55,
             size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, name, x+0.01, y+kh*0.68, kw-0.02, kh*0.28,
             size=3.2, color=LGREY, align=PP_ALIGN.CENTER)
    if badge_color and badge_txt:
        bw = 0.09
        add_rect(slide, x+kw-bw-0.01, y+0.01, bw, 0.07, fill=badge_color)
        add_text(slide, badge_txt, x+kw-bw-0.01, y+0.005, bw, 0.07,
                 size=2.8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def small_phone_reminder(slide, l=11.1, t=1.9, label="Data Source: BiziBid User Data"):
    """Small reminder phone in corner with data-source label."""
    w = 1.75
    sl, st, sw, sh = draw_phone(slide, l, t, w=w)
    ph_fill(slide, sl, st, sw, sh, color=APP_DARK)
    y = ph_status(slide, sl, st, sw)
    y = ph_header(slide, sl, y, sw, "BiziBid", bg=GOLD, fg=BLACK)
    add_rect(slide, sl, y, sw, sh - (y - st), fill=APP_DARK)
    kw = (sw - 0.015) / 2
    kh = kw * 0.88
    stores = ["Massy", "Carters", "H&B", "Price-R"]
    badges = [(RED, "3"), None, (GREEN, "1"), None]
    gap = 0.012
    for r in range(2):
        for c in range(2):
            i = r * 2 + c
            bc, bt = badges[i] if badges[i] else (None, "")
            ph_store_key(slide, sl + c*(kw+gap), y + gap + r*(kh+gap), kw, kh,
                         stores[i], badge_color=bc, badge_txt=bt)
    phone_h = w * 2.08
    add_text(slide, label, l-0.05, t + phone_h + 0.07, w+0.1, 0.2,
             size=6, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 5.2, fill=NAVY)
add_rect(s, 0, 5.2, 13.33, 0.08, fill=GOLD)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.1), Inches(11), Inches(1.8))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Bizi"; r1.font.size = Pt(80); r1.font.bold = True; r1.font.color.rgb = WHITE
r2 = p.add_run(); r2.text = "Bid";  r2.font.size = Pt(80); r2.font.bold = True; r2.font.color.rgb = GOLD
add_text(s, "Barbados' First Retail Auction & Deals Platform",
         1.0, 3.05, 11.3, 0.6, size=20, color=RGBColor(0xC0,0xC8,0xD8), align=PP_ALIGN.CENTER)
add_text(s, "Partnership Opportunity  ·  Confidential",
         1.0, 3.75, 11.3, 0.5, size=13, color=RGBColor(0x5A,0x66,0x78),
         align=PP_ALIGN.CENTER, italic=True)
add_text(s, "Presented by INTERXDB  ·  Digital Retail & Auction Platform  ·  Barbados",
         1.0, 5.5, 11.3, 0.45, size=11, color=MGREY, align=PP_ALIGN.CENTER)
add_text(s, "bizibid.com  ·  contact@interxdb.com  ·  Tel 1(246) 241-3771",
         1.0, 6.0, 11.3, 0.4, size=11, color=LGREY, align=PP_ALIGN.CENTER)
page_num(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS BIZIBID?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Platform Overview")
heading(s, "What is BiziBid?")
sub(s, "A mobile-first auction and deals platform built exclusively for Barbados' retail market.")
pillars = [
    ("★", "The BiziBid Mobile App",
     "A free app (iOS & Android) where registered customers browse live auctions and exclusive deals — 24 hours a day, 7 days a week."),
    ("★", "Your Back-Office Portal",
     "A private management portal at BiziBid.com. Upload products, create deals, review auction results, and track performance — from any browser."),
    ("★", "Fully Operated by INTERXDB",
     "INTERXDB owns and operates every part of the platform — the app, the servers, the technology. You simply list your products. We handle everything else."),
]
for i, (icon, title, body) in enumerate(pillars):
    feature_card(s, icon, title, body, l=0.5 + i*4.3, t=2.7, w=4.1, h=2.6)
page_num(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE BUZZ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, BUZZ_DARK)

# Gold accent bar near top
add_rect(s, 0, 0, 13.33, 0.06, fill=GOLD)

# Logo on dark bg
tb = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(2.5), Inches(0.55))
tf = tb.text_frame; p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "Bizi"; r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = WHITE
r2 = p.add_run(); r2.text = "Bid";  r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = GOLD

# "SOMETHING NEW IS COMING" badge
add_rect(s, 0.5, 0.88, 3.4, 0.38, fill=BUZZ_BLUE)
add_text(s, "SOMETHING NEW IS COMING TO BARBADOS", 0.6, 0.91, 3.2, 0.32,
         size=9, bold=True, color=WHITE)

# Main headline
add_text(s, "Wake the Market.", 0.5, 1.4, 12.3, 1.1,
         size=52, bold=True, color=WHITE)
add_text(s, "Your customers are ready for something exciting — are you first?",
         0.5, 2.55, 10.5, 0.65, size=18, color=RGBColor(0xC0, 0xC8, 0xD8))

# Three buzz boxes
buzz_items = [
    ("The buzz is building.", "Barbados' retail customers are craving deals, excitement, and a reason to engage with their favourite stores again. BiziBid is that reason."),
    ("First mover wins.", "Founding retail partners secure the most visible positions in the app before it opens to competitors. The stores that join first are the stores that dominate."),
    ("Your brand stays alive 24/7.", "Even when your store doors are closed, your name, your products, and your deals are on every customer's phone — generating excitement around the clock."),
]
for i, (bold_txt, body_txt) in enumerate(buzz_items):
    x = 0.5 + i * 4.28
    add_rect(s, x, 3.42, 4.05, 2.82, fill=RGBColor(0x0D, 0x1E, 0x40))
    add_rect(s, x, 3.42, 0.06, 2.82, fill=BUZZ_AMBER)
    add_text(s, bold_txt, x+0.16, 3.54, 3.75, 0.42,
             size=13, bold=True, color=BUZZ_AMBER)
    add_text(s, body_txt, x+0.16, 4.0, 3.75, 2.1,
             size=11, color=RGBColor(0xC0, 0xC8, 0xD8))

add_text(s, "Early partners lock in rates and placement for 12 months. Positions are limited.",
         0.5, 6.45, 12.3, 0.38, size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
page_num(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — USER REGISTRATION & DEMOGRAPHICS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Customer Intelligence")
heading(s, "Know Your Customers Before They Walk In")
sub(s, "Every BiziBid user registers with verified personal details — giving your store valuable demographic intelligence.")

# Left content
bullet_box(s, [
    "★  Full name, verified email, and mobile number collected at registration",
    "★  Date of birth — every user must be 18 or older — you know your age demographic",
    "★  Parish / address — you see exactly which areas your customers come from",
    "★  All data is aggregated and visible in your back-office statistics dashboard",
    "★  Track which age groups bid most, which parishes are most active, and more",
    "★  No other advertising platform in Barbados gives retail stores this level of insight",
], l=0.5, t=2.55, w=9.55, h=4.3, text_color=BLACK, text_size=13)

add_rect(s, 0.5, 6.4, 9.55, 0.72, fill=LIGHT)
add_rect(s, 0.5, 6.4, 0.06, 0.72, fill=GOLD)
add_text(s, "Every registered user is a real, verified person — not an anonymous click. Your store sees the name, age group, and parish of every customer who follows, bids on, or likes your products.",
         0.65, 6.46, 9.3, 0.62, size=10, italic=True, color=MGREY)

# Phone mockup — registration screen
sl, st, sw, sh = draw_phone(s, l=10.35, t=1.3, w=2.52)
ph_fill(s, sl, st, sw, sh, color=APP_DARK)
y = ph_status(s, sl, st, sw)
y = ph_header(s, sl, y, sw, "Create Account")
y += 0.04
y = ph_label(s, sl, y, sw, "Tell us about yourself", color=LGREY)
y = ph_field(s, sl, y, sw, "First Name", "John")
y = ph_field(s, sl, y, sw, "Last Name", "Husbands")
y = ph_field(s, sl, y, sw, "Date of Birth (18+ required)", "15 / 04 / 1992")
y = ph_field(s, sl, y, sw, "Parish", "St. Michael")
y = ph_field(s, sl, y, sw, "Email Address", "john@gmail.com")
y = ph_field(s, sl, y, sw, "Mobile Number", "+1 (246) 256-8821")
y += 0.06
ph_btn(s, sl, y, sw, "REGISTER FREE")

# Age / parish callout labels next to phone
add_text(s, "← Age verified (18+)",    10.95, 3.05, 2.2, 0.3, size=8, color=BUZZ_BLUE, bold=True)
add_text(s, "← Parish captured",       10.95, 3.82, 2.2, 0.3, size=8, color=BUZZ_BLUE, bold=True)
add_text(s, "← Real contact info",     10.95, 4.58, 2.2, 0.3, size=8, color=BUZZ_BLUE, bold=True)
page_num(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INSIDE THE APP (large phone, home screen)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "The BiziBid App")
heading(s, "Inside the App — Your Store on Every Screen")
sub(s, "The BiziBid home screen is a live keypad of retail stores — your logo always visible, always active.")

# Feature callouts — left column
callouts = [
    ("PINNED STORE", "Customers pin their favourite store at the top — maximum visibility every time they open the app."),
    ("LIVE BADGE", "Colour-coded badge shows active bids, deals, and alerts on your store button — customers react immediately."),
    ("STORE KEYPAD", "Your logo is a direct-tap button. One press opens your full auction and deals screen instantly."),
    ("DEAL STRIP", "Active Deals of the Week and Month appear as a scrollable banner across the top of the feed."),
]
for i, (lbl, desc) in enumerate(callouts):
    yt = 2.55 + i * 1.17
    add_rect(s, 0.4, yt, 6.7, 1.05, fill=LIGHT)
    add_rect(s, 0.4, yt, 0.05, 1.05, fill=GOLD)
    add_text(s, lbl,  0.55, yt+0.07, 6.4, 0.26, size=9, bold=True, color=GOLD)
    add_text(s, desc, 0.55, yt+0.34, 6.4, 0.64, size=10, color=MGREY)

# Large phone (center-right)
sl, st, sw, sh = draw_phone(s, l=7.45, t=0.72, w=3.05)
ph_fill(s, sl, st, sw, sh, color=APP_DARK)
y = ph_status(s, sl, st, sw)
y = ph_header(s, sl, y, sw, "BiziBid")

# Deal strip at top of feed
add_rect(s, sl, y, sw, 0.18, fill=RGBColor(0x0A, 0x1A, 0x30))
add_text(s, "DEAL OF THE WEEK: 32\" Smart TV  $199", sl+0.03, y+0.03, sw-0.05, 0.12,
         size=3.8, bold=True, color=GOLD)
y += 0.18

y = ph_label(s, sl, y+0.02, sw, "YOUR STORES", color=LGREY)

# Keypad grid — 3 columns, 2 rows
kw = (sw - 0.025) / 3
kh = kw * 0.9
gap = 0.01
stores_kp = [
    ("Massy",    RED,   "5"),
    ("Carters",  GREEN, "2"),
    ("H&B",      None,  ""),
    ("Price-R",  RED,   "1"),
    ("Courts",   None,  ""),
    ("MasonMrt", GREEN, "3"),
]
for idx, (name, bc, bt) in enumerate(stores_kp):
    row, col = idx // 3, idx % 3
    ph_store_key(s,
                 sl + col*(kw+gap), y + gap + row*(kh+gap),
                 kw, kh, name, badge_color=bc, badge_txt=bt)

y += 2*(kh+gap) + gap + 0.06
y = ph_label(s, sl, y, sw, "MY ACTIVE BIDS", color=LGREY)
y = ph_bid_row(s, sl, y, sw, "32\" Smart TV", "$215", winning=True)
y = ph_bid_row(s, sl, y, sw, "Gas BBQ Grill", "$88")

# Arrows pointing from callouts to phone
add_text(s, "→", 7.1, 2.92, 0.3, 0.4, size=14, bold=True, color=GOLD)
add_text(s, "→", 7.1, 3.62, 0.3, 0.4, size=14, bold=True, color=GOLD)
add_text(s, "→", 7.1, 4.32, 0.3, 0.4, size=14, bold=True, color=GOLD)

page_num(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW CUSTOMERS BID
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "The Bidding Experience")
heading(s, "How Customers Bid — Simple, Fast, Addictive")
sub(s, "Three taps and your customer is competing live. The simplicity is what drives repeat engagement.")

# Steps — left side
steps = [
    ("1", "Open BiziBid",       "Customer opens the app and sees your store on the home screen with an active bid badge."),
    ("2", "Tap Your Store",     "One tap opens your auction listing — product image, description, current bid price, and live countdown."),
    ("3", "Tap BID NOW",        "A single tap places their bid. They are instantly in the lead. No forms, no complicated process."),
    ("4", "Outbid Alert Fires", "If outbid, a push notification arrives immediately on their phone bringing them straight back."),
    ("5", "Win & Collect",      "Auction closes, winner receives a unique collection code. They bring it to your store and collect their item."),
]
for i, (num, title, body) in enumerate(steps):
    yt = 2.55 + i * 0.97
    add_rect(s, 0.4, yt, 0.42, 0.42, fill=NAVY)
    add_text(s, num, 0.4, yt+0.04, 0.42, 0.35, size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, title, 0.92, yt+0.01, 3.9, 0.28, size=12, bold=True, color=BLACK)
    add_text(s, body,  0.92, yt+0.30, 3.9, 0.60, size=10, color=MGREY)

add_rect(s, 4.88, 2.55, 5.0, 4.82, fill=LIGHT)
add_rect(s, 4.88, 2.55, 0.06, 4.82, fill=GOLD)
add_text(s, "Bidding is designed to feel like a game — customers come back to defend their position.",
         5.05, 2.65, 4.65, 0.55, size=11, italic=True, color=MGREY)
add_text(s, "No credit card needed. No online payment. Winners collect and pay cash in your store.",
         5.05, 3.25, 4.65, 0.55, size=11, italic=True, color=MGREY)
add_text(s, "Every bid placed is a moment your store name is top-of-mind for that customer.",
         5.05, 3.85, 4.65, 0.55, size=11, italic=True, color=MGREY)
add_text(s, "Auctions run even while your store is closed — you wake up to completed sales.",
         5.05, 4.45, 4.65, 0.55, size=11, italic=True, color=MGREY)
add_text(s, "Multiple stores. Multiple auctions. One app. Customers manage all bids in one place.",
         5.05, 5.05, 4.65, 0.55, size=11, italic=True, color=MGREY)
add_text(s, "Each auction winner is stored in your portal with their name and collection code.",
         5.05, 5.65, 4.65, 0.55, size=11, italic=True, color=MGREY)

# Phone — auction detail screen
sl, st, sw, sh = draw_phone(s, l=10.35, t=1.3, w=2.52)
ph_fill(s, sl, st, sw, sh, color=APP_DARK)
y = ph_status(s, sl, st, sw)
y = ph_header(s, sl, y, sw, "Live Auction")

# Product display area
add_rect(s, sl, y, sw, 0.65, fill=RGBColor(0x08, 0x18, 0x30))
add_text(s, "32\" Smart TV — Samsung", sl+0.05, y+0.04, sw-0.08, 0.22, size=5.5, bold=True, color=WHITE)
add_text(s, "Starting bid: $150.00", sl+0.05, y+0.26, sw-0.08, 0.14, size=4.5, color=LGREY)
add_text(s, "Ends in  00:04:32", sl+0.05, y+0.43, sw-0.08, 0.17, size=4.8, bold=True, color=RED)
y += 0.65

# Current bid banner
add_rect(s, sl, y, sw, 0.24, fill=NAVY)
add_text(s, "CURRENT BID", sl+0.04, y+0.02, sw*0.5, 0.09, size=3.8, color=LGREY)
add_text(s, "$215.00", sl+0.04, y+0.11, sw*0.6, 0.12, size=9, bold=True, color=GOLD)
y += 0.24

y = ph_label(s, sl, y+0.04, sw, "BID HISTORY", color=LGREY)
y = ph_bid_row(s, sl, y, sw, "j.husbands",  "$215.00", winning=True)
y = ph_bid_row(s, sl, y, sw, "r.alleyne",   "$200.00")
y = ph_bid_row(s, sl, y, sw, "b.springer",  "$185.00")

y += 0.06
ph_btn(s, sl, y, sw, "BID NOW — $220.00", bg=GOLD, fg=BLACK)
page_num(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BID & OUTBID NOTIFICATIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Real-Time Notifications")
heading(s, "Notifications That Bring Customers Back Instantly")
sub(s, "Every bid event triggers a push notification — customers never miss a moment of the auction.")

# Left: two notification types explained
add_rect(s, 0.4, 2.55, 9.65, 2.05, fill=RGBColor(0xF0, 0xFB, 0xF4))
add_rect(s, 0.4, 2.55, 0.06, 2.05, fill=GREEN)
add_text(s, "WINNING BID NOTIFICATION", 0.56, 2.63, 9.2, 0.32, size=10, bold=True, color=GREEN)
add_text(s, "Sent immediately when a customer takes the lead. The positive reinforcement is powerful — it keeps them emotionally invested and returning to the app to check their status.",
         0.56, 3.0, 9.2, 0.55, size=12, color=MGREY)
add_text(s, "Result: Customer stays engaged, monitors the auction, watches rivals.",
         0.56, 3.6, 9.2, 0.4, size=11, bold=True, color=BLACK)
add_text(s, "  'You're winning! Keep it up.'", 0.56, 4.05, 9.2, 0.4, size=12, italic=True, color=MGREY)

add_rect(s, 0.4, 4.74, 9.65, 2.1, fill=RGBColor(0xFD, 0xF2, 0xF2))
add_rect(s, 0.4, 4.74, 0.06, 2.1, fill=RED)
add_text(s, "OUTBID NOTIFICATION", 0.56, 4.82, 9.2, 0.32, size=10, bold=True, color=RED)
add_text(s, "Fires the instant another customer places a higher bid. The urgency is immediate — customers tap straight back into the app to respond. Most counter-bids happen within 90 seconds.",
         0.56, 5.19, 9.2, 0.55, size=12, color=MGREY)
add_text(s, "Result: Customer returns to app, counter-bids, and engagement spikes.",
         0.56, 5.79, 9.2, 0.4, size=11, bold=True, color=BLACK)
add_text(s, "  'You've been outbid! Tap to stay in the lead.'", 0.56, 6.24, 9.2, 0.4, size=12, italic=True, color=MGREY)

# Phone — notification screen
sl, st, sw, sh = draw_phone(s, l=10.35, t=1.3, w=2.52)
ph_fill(s, sl, st, sw, sh, color=APP_DARK)
y = ph_status(s, sl, st, sw)
y = ph_header(s, sl, y, sw, "Notifications")

y = ph_notif(s, sl, y, sw,
             "You're winning!",
             "32\" Smart TV — bid $215",
             badge=NOTIF_GRN, ts="now")
y = ph_notif(s, sl, y, sw,
             "Your bid is leading",
             "Gas BBQ Grill — bid $88",
             badge=NOTIF_GRN, ts="2 min")
y += 0.015
y = ph_notif(s, sl, y, sw,
             "You've been outbid!",
             "32\" Smart TV — bid $225",
             badge=NOTIF_RED, ts="5 min")
y = ph_notif(s, sl, y, sw,
             "Tap to counter-bid",
             "Someone bid $225 — you had $215",
             badge=NOTIF_RED, ts="5 min")
y += 0.015
y = ph_notif(s, sl, y, sw,
             "Auction ending soon!",
             "Gas BBQ Grill closes in 2 hrs",
             badge=RGBColor(0xD9, 0x77, 0x06), ts="1 hr")
y += 0.04
ph_btn(s, sl, y, sw, "VIEW AUCTIONS", bg=GOLD, fg=BLACK)
page_num(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEALS OF THE WEEK & MONTH
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Feature — Deals")
heading(s, "Deals of the Week & Month — Drive Foot Traffic")
sub(s, "Fixed-price deals showcase your best offers and bring customers through your doors.")

left_items = [
    "▶  You set a fixed price — no bidding, no auction pressure",
    "▶  Customers see your deal prominently in the app all week or all month",
    "▶  Customers can Like the deal — like count is visible to all users (social proof)",
    "▶  Customers share deals directly to WhatsApp, Instagram, and other platforms",
    "▶  Shared links show the deal to people who do not yet have the app — driving new downloads",
    "▶  Walk-in customers ask for the deal by name — increased foot traffic guaranteed",
    "▶  Deals of the Month offer extended 30-day visibility at premium placement",
]
bullet_box(s, left_items, l=0.5, t=2.55, w=7.5, h=4.2, text_color=BLACK, text_size=13)

add_rect(s, 8.3, 2.55, 4.5, 1.9, fill=LIGHT)
add_rect(s, 8.3, 2.55, 0.06, 1.9, fill=GOLD)
add_text(s, "Deal of the Week",    8.5, 2.65, 4.1, 0.38, size=12, bold=True, color=BLACK)
add_text(s, "7-day fixed-price feature\nPremium placement in app\nLike + Share enabled\nDrives walk-in purchase",
         8.5, 3.1, 4.0, 1.25, size=11, color=MGREY)

add_rect(s, 8.3, 4.6, 4.5, 1.9, fill=NAVY)
add_rect(s, 8.3, 4.6, 0.06, 1.9, fill=GOLD)
add_text(s, "Deal of the Month",   8.5, 4.7,  4.1, 0.38, size=12, bold=True, color=GOLD)
add_text(s, "Full 30-day visibility\nTop placement position\nMaximum brand exposure\nIdeal for key product lines",
         8.5, 5.15, 4.0, 1.25, size=11, color=RGBColor(0xC0,0xC8,0xD8))
page_num(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — NOT JUST VIRTUAL RETAIL — A 24/7 BRAND PORTAL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Brand Strategy")
heading(s, "More Than a Marketplace — A 24/7 Brand Portal", w=10.4, size=28)
sub(s, "BiziBid is your store's permanent digital home. It keeps your brand alive in customers' minds every single day.",
    w=10.4)

portal_points = [
    ("Open 24/7",
     "Your store is visible and active on BiziBid even when your doors are closed. Customers browse your products at midnight, on weekends, on public holidays."),
    ("Not just a sale — an experience",
     "Customers don't just buy from BiziBid. They follow your store, they watch your auctions, they share your deals, and they talk about you. That is brand building."),
    ("A daily digital billboard",
     "Your store logo sits on every registered customer's phone home screen. No billboard, no newspaper ad, no social post delivers that kind of daily reach at this price."),
    ("Customer loyalty, digitally",
     "When a customer pins your store, sets a notification, and bids week after week — that is a loyal customer. BiziBid builds that relationship for you automatically."),
    ("Your competitors are still offline",
     "Most Barbados retail stores have no meaningful digital presence. BiziBid makes your store the digital leader in your category before your rivals even react."),
    ("The app is your voice",
     "Every auction, every deal, every share is your brand speaking to customers. BiziBid gives your brand a voice that never goes silent."),
]
for i, (title, body) in enumerate(portal_points):
    c, r = i % 2, i // 2
    x = 0.4 + c * 5.05
    y = 2.55 + r * 1.55
    add_rect(s, x, y, 4.8, 1.42, fill=LIGHT)
    add_rect(s, x, y, 0.05, 1.42, fill=GOLD)
    add_text(s, title, x+0.14, y+0.1, 4.5, 0.32, size=11, bold=True, color=NAVY)
    add_text(s, body,  x+0.14, y+0.44, 4.5, 0.9,  size=10, color=MGREY)

# Small phone reminder — right side
small_phone_reminder(s, l=10.75, t=2.6, label="Your brand on every customer's phone")
page_num(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LIVE AUCTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Feature — Auctions")
heading(s, "Live Auctions — Excitement That Sells")
sub(s, "Customers compete in real-time for your products. The excitement keeps them engaged and coming back.")

bullet_box(s, [
    "▶  You upload products with a starting bid price (set below shelf price)",
    "▶  Up to 25 items per week rotate through the platform automatically",
    "▶  5 new items are released each day at varied morning hours — customers check in daily",
    "▶  Auctions run for 3, 7, or 14 days depending on the item",
    "▶  Real-time bidding — customers see bids update live on their screen",
    "▶  Unsold items re-enter the rotation for the following week automatically",
], l=0.5, t=2.55, w=6.0, h=4.2, text_color=BLACK, text_size=13)

add_rect(s, 7.0, 2.55, 5.8, 4.2, fill=NAVY)
add_rect(s, 7.0, 2.55, 5.8, 0.06, fill=GOLD)
add_text(s, "When a Customer Wins", 7.2, 2.65, 5.4, 0.45, size=13, bold=True, color=GOLD)
win_steps = [
    "1  Auction closes — winner is confirmed",
    "2  A unique collection code is generated",
    "3  Winner receives code on their phone",
    "4  Winner visits your store and shows code to cashier",
    "5  Cashier confirms the code — item is collected",
    "6  Record of the sale is stored in your portal",
]
bullet_box(s, win_steps, l=7.2, t=3.2, w=5.3, h=3.3, text_color=WHITE, text_size=12)
add_text(s, "No online payment needed — all purchases are cash in-store.",
         7.2, 6.35, 5.3, 0.38, size=10, italic=True, color=LGREY)
page_num(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BACK-OFFICE PORTAL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Your Back-Office — BiziBid.com")
heading(s, "Your Private Management Portal", w=10.3)
sub(s, "Log in at BiziBid.com to manage your entire BiziBid presence — available 24/7 from any browser.", w=10.3)

portal_features = [
    ("★", "Bulk Product Upload",    "Upload multiple products at once with images, descriptions, and your starting bid price. No technical knowledge required."),
    ("★", "Deals Management",       "Create and schedule Deals of the Week and Deals of the Month. Set your own price and valid dates. Activate with one click."),
    ("★", "Auction Results",        "View all completed auctions — winning bid amount, customer name, collection code. Full record of every sale."),
    ("★", "Collection Codes",       "See the unique code for each auction winner so your cashier can verify the customer quickly and confidently."),
    ("★", "Invoice Download",       "View and download your monthly platform invoices as PDF. Full record of your subscription and activity charges."),
    ("★", "Store Statistics",       "Live dashboard showing user engagement, bid activity, deal shares, and more. Know exactly how customers interact with your store."),
]
for i, (icon, title, body) in enumerate(portal_features):
    c, r = i % 3, i // 3
    feature_card(s, icon, title, body, l=0.35 + c*3.42, t=2.55 + r*1.92, w=3.22, h=1.78)

small_phone_reminder(s, l=10.78, t=2.6, label="Powered by real app user data")
page_num(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — STORE STATISTICS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Back-Office — Statistics Dashboard")
heading(s, "Know Your Customers — Live Store Statistics")
sub(s, "Your back-office portal shows you exactly how registered BiziBid users interact with your store in real time.")

stats1 = [("1,240", "Registered users\nfollowing your store"),("86","Active bids on\nyour items today"),("312","Users watching\nyour products"),("94%","Auction completion\nrate this month")]
for i,(num,lbl) in enumerate(stats1):
    stat_card(s, num, lbl, l=0.5+i*3.1, t=2.55, w=2.9, h=1.2)

stats2 = [("48","Deals shared on\nWhatsApp & Instagram"),("$8,400","Total auction value\ngenerated this month"),("22","New app users from\nyour shared links"),("4.8★","Average customer\nengagement score")]
for i,(num,lbl) in enumerate(stats2):
    stat_card(s, num, lbl, l=0.5+i*3.1, t=3.9, w=2.9, h=1.2)

add_rect(s, 0.5, 5.28, 12.3, 1.38, fill=LIGHT)
add_rect(s, 0.5, 5.28, 0.06, 1.38, fill=GOLD)
add_text(s, "All statistics are updated in real time in your back-office portal at BiziBid.com. "
            "Track which products generate the most interest, which deals are most shared, and how your store's following "
            "grows week over week. See age group breakdowns and which parishes your most active bidders come from. "
            "No data is shared with other retail partners.",
         0.65, 5.35, 12.0, 1.22, size=12, color=MGREY)
page_num(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SOCIAL SHARING: 10x TO 100x VISIBILITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 0.06, fill=GOLD)

# Logo on white
logo_text(s)
section_label(s, "Viral Sharing Power")
heading(s, "One Share = 10x to 100x Your Visibility", size=28)
sub(s, "When a customer shares your product on WhatsApp or social media, your brand reaches their entire network instantly.")

# Big number boxes
share_data = [
    ("10x",   "WhatsApp Message",    "Sent to a group of 50-150 people. Everyone in that group sees your store name, product, and price."),
    ("50x",   "WhatsApp Status",     "Visible to all of a customer's WhatsApp contacts for 24 hours. Average reach: 200-400 people per post."),
    ("100x",  "Facebook / Instagram","Liked, shared, or reposted content reaches friends-of-friends. A single viral share can reach thousands."),
]
for i, (mult, channel, body) in enumerate(share_data):
    x = 0.4 + i * 4.28
    add_rect(s, x, 2.55, 4.05, 3.2, fill=NAVY)
    add_rect(s, x, 2.55, 4.05, 0.05, fill=GOLD)
    add_text(s, mult,    x+0.15, 2.65, 3.75, 0.88, size=52, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, channel, x+0.15, 3.62, 3.75, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, body,    x+0.15, 4.08, 3.75, 1.5,  size=11, color=RGBColor(0xC0,0xC8,0xD8), align=PP_ALIGN.CENTER)

# Bottom statement
add_rect(s, 0.4, 5.9, 12.5, 1.26, fill=LIGHT)
add_rect(s, 0.4, 5.9, 0.06, 1.26, fill=GOLD)
add_text(s, "People love to share a great deal.", 0.6, 5.98, 12.1, 0.38, size=14, bold=True, color=BLACK)
add_text(s, "A customer who finds a $50 item auction starting at $12 will share that to every WhatsApp group they have. "
            "That is free marketing to hundreds of people who have never heard of your store — and it costs you nothing. "
            "Every shared link drives new app downloads, and new app users are new bidders on your products.",
         0.6, 6.4, 12.1, 0.68, size=11, color=MGREY)
page_num(s, 13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — VIRAL GROWTH FUNNEL
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Feature — Viral Growth")
heading(s, "Every Share is Free Marketing for Your Store")
sub(s, "When customers share your auctions and deals, your store name reaches people who have never heard of BiziBid.")

flow = [
    ("Your customer\nsees a great deal\nor live auction", "01"),
    ("They share it\nto WhatsApp or\nInstagram", "02"),
    ("Their contact sees\nyour store name,\nprice & product", "03"),
    ("They visit\nBiziBid.com and\ndownload the app", "04"),
    ("New registered\nbidder — on your\nauctions", "05"),
]
for i, (text, num) in enumerate(flow):
    x = 0.35 + i * 2.56
    add_rect(s, x, 2.5, 2.3, 2.3, fill=NAVY if i % 2 == 0 else LIGHT)
    add_text(s, num,  x+0.1, 2.58, 0.5, 0.45, size=11, bold=True, color=GOLD)
    add_text(s, text, x+0.1, 2.95, 2.1, 1.7, size=12, color=WHITE if i % 2 == 0 else BLACK)
    if i < 4:
        add_text(s, "->", x+2.32, 3.35, 0.22, 0.5, size=18, color=GOLD)

add_rect(s, 0.5, 5.05, 12.3, 1.65, fill=LIGHT)
add_rect(s, 0.5, 5.05, 0.06, 1.65, fill=GOLD)
bullet_box(s, [
    "Auction shares show the live bid price and countdown — creating instant urgency for the viewer.",
    "Deal shares show your store name, product image, and fixed price — driving walk-in intent.",
    "All shared links direct visitors to BiziBid.com where they can download the app. New downloads = new bidders.",
], l=0.7, t=5.12, w=12.0, h=1.5, text_color=MGREY, text_size=11)
page_num(s, 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — KEYPAD PLACEMENT & TIERS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Placement & Visibility")
heading(s, "Secure Your Position on the BiziBid Keypad", w=10.4)
sub(s, "Your store logo placement on the app home screen is a paid position. First come, first placed.", w=10.4)

tiers = [
    (GOLD,                      "Tier 1 — Pinned Position",
     "Top of every user's home screen. Your logo is displayed largest, with full notification badge visibility. One store only — the most premium placement available."),
    (NAVY,                      "Tier 2 — Primary Keypad",
     "First 6–9 visible positions in the main keypad grid. Seen immediately when the app opens. No scrolling required. Highest traffic positions after Tier 1."),
    (MGREY,                     "Tier 3 — Slide-Out Panel",
     "Your store appears in the extended keypad panel (swipe to reveal). Still accessible and notified — ideal for entry-level platform participation."),
    (RGBColor(0x44,0x44,0x44),  "Banner Advertisements",
     "A rotating advertisement strip displayed across all screens at the top of the app. Separate from keypad placement — available to any registered partner."),
]
for i, (color, title, body) in enumerate(tiers):
    x = 0.4 + (i % 2) * 5.15
    y = 2.55 + (i // 2) * 2.18
    add_rect(s, x, y, 4.9, 1.98, fill=color)
    add_text(s, title, x+0.18, y+0.14, 4.55, 0.45, size=13, bold=True,
             color=BLACK if color == GOLD else WHITE)
    add_text(s, body,  x+0.18, y+0.62, 4.55, 1.24, size=11,
             color=NAVY if color == GOLD else RGBColor(0xD0,0xD8,0xE8) if color == NAVY else WHITE)

small_phone_reminder(s, l=10.78, t=2.62, label="Your logo on every user's home screen")
page_num(s, 15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — PRICING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Partnership Pricing")
heading(s, "Simple, Transparent Pricing (BBD)")
sub(s, "All rates are fixed for the first 12 months for founding partners. No hidden fees.")

add_rect(s, 0.5, 2.4, 12.3, 1.42, fill=NAVY)
add_rect(s, 0.5, 2.4, 0.06, 1.42, fill=GOLD)
add_text(s, "60 DAYS FREE", 0.7, 2.48, 2.5, 0.52, size=22, bold=True, color=GOLD)
add_text(s, "No monthly cost for the first 60 days. A $500.00 BBD deposit secures your primary advertisement "
            "placement within the app — guaranteeing your store is featured prominently to all registered "
            "users from day one.",
         3.3, 2.5, 9.3, 1.22, size=12, color=WHITE)

headers = ["Service", "Description", "Rate (BBD)"]
rows = [
    ("Base Listing Package", "Up to 25 items listed per week on the app",            "$450.00 / month"),
    ("Additional Items",     "Per additional 10 items above the 25-item base",        "$150.00 / month"),
    ("Deal of the Week",     "Featured deal highlighted to all users for 7 days",     "$150.00 / week"),
    ("Deal of the Month",    "Premium featured deal — top placement for full month",  "$200.00 / month"),
]
th = 0.38; rh = 0.52; table_t = 3.98
col_w = [3.8, 6.2, 2.1]; col_x = [0.5, 4.3, 10.5]

for j,(hdr,cx,cw) in enumerate(zip(headers,col_x,col_w)):
    add_rect(s, cx, table_t, cw, th, fill=NAVY)
    add_text(s, hdr, cx+0.1, table_t+0.05, cw-0.15, th-0.08, size=11, bold=True, color=WHITE)

for i, row in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    for j,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
        add_rect(s, cx, table_t+th+i*rh, cw, rh, fill=bg)
        add_text(s, cell, cx+0.1, table_t+th+i*rh+0.07, cw-0.15, rh-0.1,
                 size=11, bold=(j==2), color=NAVY if j==2 else BLACK)

add_text(s, "All pricing in Barbados Dollars (BBD). Rates fixed for the first 12 months for founding partners. "
            "Keypad placement and banner ad pricing available on request.",
         0.5, 7.0, 12.3, 0.38, size=10, italic=True, color=LGREY)
page_num(s, 16)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — WHY BIZIBID
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE); gold_bar(s); logo_text(s)
section_label(s, "Why BiziBid?")
heading(s, "Built for Barbados Retail. Starting Today.")

reasons = [
    ("Barbados First",          "BiziBid is built specifically for the Barbados market — the culture, the community, the retail landscape. Not a foreign platform adapted for the island."),
    ("Zero Technical Effort",   "You upload products and set prices. INTERXDB operates the entire platform — the app, the hosting, the technology. Nothing more is required from you."),
    ("Keep Your Store Relevant","Digital shopping is growing. BiziBid gives your store a daily digital presence that keeps you competitive — without the cost of building your own platform."),
    ("Exclusive & Professional","Only verified, registered businesses may list products. No individual sellers. No flea-market feel. A premium marketplace that reflects well on your brand."),
    ("Your Target Market",      "BiziBid targets medium-to-high income earners and family units — the customers who already shop at your store. Reach them digitally where they are."),
    ("Founding Partner Advantage","Founding partners lock in rates for 12 months and secure preferred keypad placement before the platform opens to all retailers. Act now — positions are limited."),
]
for i, (title, body) in enumerate(reasons):
    c, r = i % 3, i // 3
    feature_card(s, "★", title, body, l=0.4 + c*4.3, t=2.55 + r*2.0, w=4.05, h=1.82)
page_num(s, 17)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 5.5, fill=NAVY)
add_rect(s, 0, 5.5, 13.33, 0.08, fill=GOLD)
add_text(s, "Ready to Join BiziBid?",
         0.8, 0.7, 11.7, 0.9, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Secure your founding partner status today — positions are limited.",
         0.8, 1.65, 11.7, 0.55, size=16, color=LGREY, align=PP_ALIGN.CENTER)

actions = [
    ("Phone", "Schedule a Presentation",   "We come to you. A personalised walkthrough of the platform tailored to your store — at your convenience."),
    ("Web",   "View Your Interactive Demo", "A personalised before-and-after demonstration of how BiziBid would represent your store has been prepared for you."),
    ("Email", "Contact Us Directly",        "Our team is ready to answer any questions and walk you through the partnership terms at your pace."),
]
for i, (icon, title, body) in enumerate(actions):
    x = 0.8 + i * 3.95
    add_rect(s, x, 2.55, 3.65, 2.65, fill=RGBColor(0x12,0x2A,0x4A))
    add_rect(s, x, 2.55, 3.65, 0.05, fill=GOLD)
    add_text(s, icon[0], x+0.18, 2.65, 0.55, 0.5, size=22, color=GOLD)
    add_text(s, title,   x+0.18, 3.22, 3.3,  0.5, size=12, bold=True, color=WHITE)
    add_text(s, body,    x+0.18, 3.72, 3.3,  1.3, size=10, color=LGREY)

tb = s.shapes.add_textbox(Inches(0.8), Inches(5.78), Inches(11.7), Inches(0.55))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Bizi"; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = BLACK
r2 = p.add_run(); r2.text = "Bid";  r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = GOLD
r3 = p.add_run(); r3.text = "  ·  Operated by INTERXDB"
r3.font.size = Pt(12); r3.font.color.rgb = MGREY
add_text(s, "contact@interxdb.com   ·   Tel 1(246) 241-3771   ·   BiziBid.com",
         0.8, 6.42, 11.7, 0.4, size=12, color=MGREY, align=PP_ALIGN.CENTER)
add_text(s, "Confidential — For Addressee Only",
         0.8, 6.9, 11.7, 0.35, size=10, italic=True, color=LGREY, align=PP_ALIGN.CENTER)
page_num(s, 18)


# ── SAVE PPTX ────────────────────────────────────────────────────────────────
out_dir = r"C:\AI Software\websites\Marketing"
pptx_path = os.path.join(out_dir, "BiziBid-Partner-Presentation.pptx")
prs.save(pptx_path)
print(f"Saved: {pptx_path}")


# ══════════════════════════════════════════════════════════════════════════════
# HTML PRINT VERSION
# ══════════════════════════════════════════════════════════════════════════════

css = """
@page { size: 13.33in 7.5in; margin: 0; }
* { box-sizing: border-box; margin: 0; padding: 0; }
body { font-family: 'Segoe UI', Arial, sans-serif; background: #ccc; }
.slide {
  width: 13.33in; height: 7.5in; background: #fff;
  position: relative; page-break-after: always; overflow: hidden;
  display: flex; flex-direction: column; padding: 0.45in 0.5in 0.4in;
}
.slide:last-child { page-break-after: avoid; }
.gold-bar { position: absolute; bottom: 0; left: 0; right: 0; height: 0.09in; background: #C9A84C; }
.logo-top { position: absolute; top: 0.2in; left: 0.4in; font-size: 22pt; font-weight: 900; }
.logo-top .bizi { color: #0A0A0A; }
.logo-top .bid  { color: #C9A84C; }
.page-num { position: absolute; bottom: 0.2in; right: 0.4in; font-size: 9pt; color: #9CA3AF; }
.slide-label { font-size: 9pt; font-weight: 700; text-transform: uppercase; letter-spacing: 2px; color: #C9A84C; margin-top: 0.55in; margin-bottom: 0.05in; }
.slide-h1 { font-size: 28pt; font-weight: 900; color: #0A0A0A; line-height: 1.15; margin-bottom: 0.07in; }
.slide-sub { font-size: 13pt; color: #4B5563; margin-bottom: 0.18in; }

.card-row { display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.14in; }
.card-row.two { grid-template-columns: repeat(2, 1fr); }
.card-row.four { grid-template-columns: repeat(4, 1fr); }
.card { background: #F8F8F8; padding: 0.16in; border-left: 0.05in solid #C9A84C; }
.card.sm { padding: 0.11in; }
.card-icon { font-size: 18pt; margin-bottom: 0.04in; }
.card-title { font-size: 11pt; font-weight: 700; color: #0A0A0A; margin-bottom: 0.05in; }
.card p, .card-body { font-size: 10pt; color: #4B5563; line-height: 1.45; }
.card.sm p { font-size: 9pt; }

.big-bullets { list-style: none; padding: 0; }
.big-bullets li { font-size: 13pt; color: #0A0A0A; padding: 0.05in 0 0.05in 0.2in; position: relative; border-bottom: 1px solid #F3F4F6; }
.big-bullets li::before { content: "▸"; color: #C9A84C; position: absolute; left: 0; }

.two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 0.2in; }
.dark-box { background: #0B1F3A; padding: 0.18in; }
.dark-box-title { font-size: 12pt; font-weight: 700; color: #C9A84C; margin-bottom: 0.09in; }
.win-steps { padding-left: 0.18in; }
.win-steps li { font-size: 11pt; color: #fff; padding: 0.04in 0; }
.dark-note { font-size: 9pt; color: #9CA3AF; margin-top: 0.09in; font-style: italic; }

.stat-row { display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.12in; }
.stat-card { background: #F8F8F8; padding: 0.12in; text-align: center; border-top: 0.04in solid #C9A84C; }
.stat-num { font-size: 26pt; font-weight: 900; color: #C9A84C; }
.stat-lbl { font-size: 9pt; color: #4B5563; margin-top: 0.04in; line-height: 1.3; }

.footnote-box { background: #F8F8F8; border-left: 0.05in solid #C9A84C; padding: 0.12in; font-size: 10pt; color: #4B5563; line-height: 1.5; margin-top: 0.14in; }
.offer-box { background: #0B1F3A; color: #fff; padding: 0.15in; border-left: 0.06in solid #C9A84C; font-size: 12pt; line-height: 1.5; margin-bottom: 0.16in; }
.offer-days { font-size: 18pt; font-weight: 900; color: #C9A84C; margin-right: 0.1in; }
.price-table { width: 100%; border-collapse: collapse; font-size: 11pt; }
.price-table thead tr { background: #0B1F3A; }
.price-table thead th { padding: 0.09in; text-align: left; color: #fff; font-size: 10pt; }
.price-table tbody td { padding: 0.08in 0.09in; color: #0A0A0A; border-bottom: 1px solid #E5E7EB; }
.price-table tbody tr.alt td { background: #F8F8F8; }
.price-table tbody td.price { font-weight: 700; color: #0B1F3A; }
.table-note { font-size: 9pt; color: #9CA3AF; margin-top: 0.09in; font-style: italic; }

.flow-row { display: flex; align-items: stretch; margin: 0.1in 0; }
.flow-box { flex: 1; padding: 0.14in; font-size: 11pt; line-height: 1.4; }
.flow-box.dark  { background: #0B1F3A; color: #fff; }
.flow-box.light { background: #F8F8F8; color: #0A0A0A; }
.flow-num { display: block; font-size: 9pt; font-weight: 700; color: #9CA3AF; margin-bottom: 0.05in; }
.flow-num.gold { color: #C9A84C; }
.flow-arrow { display: flex; align-items: center; font-size: 18pt; color: #C9A84C; padding: 0 0.05in; background: #F8F8F8; }

.tier-row { display: grid; grid-template-columns: repeat(4, 1fr); gap: 0.14in; margin-top: 0.1in; }
.tier { padding: 0.18in; }
.tier.gold-tier  { background: #C9A84C; color: #0A0A0A; }
.tier.navy-tier  { background: #0B1F3A; color: #fff; }
.tier.grey-tier  { background: #4B5563; color: #fff; }
.tier.dark-tier  { background: #1F2937; color: #fff; }
.tier-label { font-size: 11pt; font-weight: 700; margin-bottom: 0.09in; }
.tier p { font-size: 10pt; line-height: 1.45; opacity: 0.9; }
.gold-tier p { color: #2A1800; }

.cover-slide { width: 100%; height: 100%; display: flex; flex-direction: column; padding: 0; }
.cover-band { background: #0B1F3A; flex: 1; display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 0.5in; }
.cover-logo { font-size: 72pt; font-weight: 900; letter-spacing: -2px; line-height: 1; }
.cover-logo .bizi { color: #fff; }
.cover-logo .bid  { color: #C9A84C; }
.cover-tagline { font-size: 18pt; color: #C0C8D8; margin-top: 0.14in; text-align: center; }
.cover-sub { font-size: 12pt; color: #5A6678; margin-top: 0.09in; font-style: italic; text-align: center; }
.cover-footer { background: #fff; padding: 0.18in 0.4in; font-size: 11pt; color: #4B5563; text-align: center; }
.cta-cards { display: flex; gap: 0.2in; margin-top: 0.28in; }
.cta-card { background: #122A4A; padding: 0.18in; flex: 1; border-top: 0.05in solid #C9A84C; }
.cta-icon { font-size: 18pt; margin-bottom: 0.07in; }
.cta-title { font-size: 12pt; font-weight: 700; color: #fff; margin-bottom: 0.06in; }
.cta-card p { font-size: 10pt; color: #9CA3AF; }

/* Buzz slide */
.buzz-slide { background: #06122A; }
.buzz-logo { font-size: 22pt; font-weight: 900; }
.buzz-logo .bizi { color: #fff; }
.buzz-logo .bid  { color: #C9A84C; }
.buzz-badge { display: inline-block; background: #1E50D8; color: #fff; font-size: 9pt; font-weight: 700; padding: 4px 10px; margin-bottom: 0.15in; }
.buzz-h1 { font-size: 52pt; font-weight: 900; color: #fff; line-height: 1; margin-bottom: 0.1in; }
.buzz-sub { font-size: 18pt; color: #C0C8D8; margin-bottom: 0.22in; }
.buzz-boxes { display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.14in; }
.buzz-box { background: #0D1E40; padding: 0.18in; border-left: 0.06in solid #F59E0B; }
.buzz-box-title { font-size: 13pt; font-weight: 700; color: #F59E0B; margin-bottom: 0.08in; }
.buzz-box p { font-size: 11pt; color: #C0C8D8; line-height: 1.5; }
.buzz-footer { font-size: 11pt; font-weight: 700; color: #C9A84C; text-align: center; margin-top: 0.18in; }

/* Phone mockup */
.with-phone { display: flex; gap: 0.25in; align-items: flex-start; }
.with-phone .text-side { flex: 1; }
.phone-wrap { flex-shrink: 0; }
.phone-outer { width: 195px; background: #101010; border-radius: 28px; padding: 17px 9px 14px; position: relative; }
.phone-speaker { width: 52px; height: 4px; background: #4B5563; border-radius: 2px; margin: 0 auto 8px; }
.phone-screen { background: #0E2A44; border-radius: 4px; overflow: hidden; }
.ph-status { background: #0E2A44; padding: 3px 7px; display: flex; justify-content: space-between; font-size: 6.5pt; color: #fff; font-weight: 700; }
.ph-header { background: #C9A84C; padding: 5px; text-align: center; font-size: 8pt; font-weight: 700; color: #000; }
.ph-header.dark { background: #0B1F3A; color: #fff; }
.ph-field { background: #12122A; padding: 3px 7px 5px; border-bottom: 1px solid #101010; }
.ph-field-lbl { font-size: 6pt; color: #9CA3AF; }
.ph-field-val { font-size: 8pt; font-weight: 700; color: #fff; }
.ph-notif { padding: 4px 7px 5px 10px; background: #12122A; border-bottom: 1px solid #101010; border-left: 3px solid #059669; position: relative; }
.ph-notif.red { border-left-color: #C02020; }
.ph-notif.amber { border-left-color: #D97706; }
.ph-notif-time { float: right; font-size: 5.5pt; color: #9CA3AF; }
.ph-notif-title { font-size: 7pt; font-weight: 700; color: #fff; margin-bottom: 1px; }
.ph-notif-body  { font-size: 6pt; color: #9CA3AF; }
.ph-btn { margin: 6px 14px 4px; background: #C9A84C; color: #000; text-align: center; padding: 5px; font-size: 8pt; font-weight: 700; border-radius: 2px; }
.ph-label { font-size: 6pt; color: #9CA3AF; padding: 3px 7px 1px; background: #0E2A44; text-transform: uppercase; letter-spacing: 1px; }
.keypad-grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 2px; padding: 3px; background: #0E2A44; }
.keypad-grid.two-col { grid-template-columns: repeat(2, 1fr); }
.store-key { background: #12122A; aspect-ratio: 1; display: flex; flex-direction: column; align-items: center; justify-content: center; position: relative; padding: 4px 2px; }
.store-key-init { font-size: 13pt; font-weight: 700; color: #C9A84C; line-height: 1; }
.store-key-name { font-size: 5pt; color: #9CA3AF; text-align: center; margin-top: 2px; }
.store-badge { position: absolute; top: 2px; right: 2px; background: #DC2626; color: #fff; font-size: 5pt; padding: 1px 3px; }
.store-badge.green { background: #16A34A; }
.deal-strip { background: #0A1A30; padding: 3px 7px; font-size: 6.5pt; font-weight: 700; color: #C9A84C; }
.bid-row { display: flex; justify-content: space-between; padding: 4px 7px; background: #12122A; border-bottom: 1px solid #101010; font-size: 7pt; }
.bid-row.top { background: #18283A; border-left: 3px solid #16A34A; }
.bid-user { color: #9CA3AF; }
.bid-row.top .bid-user { color: #fff; }
.bid-amt  { color: #C9A84C; font-weight: 700; }
.ph-prod  { background: #081830; padding: 5px 7px; }
.ph-prod-name { font-size: 7pt; font-weight: 700; color: #fff; }
.ph-prod-sub  { font-size: 6pt; color: #9CA3AF; margin-top: 1px; }
.ph-prod-timer { font-size: 6.5pt; font-weight: 700; color: #DC2626; margin-top: 1px; }
.ph-cur-bid { background: #0B1F3A; padding: 4px 7px; }
.ph-cur-lbl { font-size: 5.5pt; color: #9CA3AF; }
.ph-cur-amt { font-size: 12pt; font-weight: 900; color: #C9A84C; }
.phone-home-bar { width: 40px; height: 3px; background: #4B5563; border-radius: 2px; margin: 7px auto 0; }

/* Callout steps */
.step-row { display: flex; gap: 0.12in; align-items: flex-start; margin-bottom: 0.1in; }
.step-num { min-width: 0.36in; height: 0.36in; background: #0B1F3A; color: #C9A84C; font-size: 16pt; font-weight: 900; display: flex; align-items: center; justify-content: center; flex-shrink: 0; }
.step-title { font-size: 12pt; font-weight: 700; color: #0A0A0A; margin-bottom: 3px; }
.step-body  { font-size: 10pt; color: #4B5563; line-height: 1.4; }

/* Share slide */
.share-boxes { display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.14in; margin-bottom: 0.16in; }
.share-box { background: #0B1F3A; padding: 0.15in; text-align: center; border-top: 0.05in solid #C9A84C; }
.share-mult { font-size: 52pt; font-weight: 900; color: #C9A84C; line-height: 1; }
.share-ch   { font-size: 13pt; font-weight: 700; color: #fff; margin: 6px 0; }
.share-body { font-size: 10pt; color: #C0C8D8; line-height: 1.45; }
.share-footer { background: #F8F8F8; border-left: 0.06in solid #C9A84C; padding: 0.13in; }
.share-footer-h { font-size: 14pt; font-weight: 700; color: #0A0A0A; margin-bottom: 0.06in; }
.share-footer p { font-size: 11pt; color: #4B5563; line-height: 1.5; }

/* Demographic callout boxes */
.demo-boxes { display: grid; grid-template-columns: repeat(3, 1fr); gap: 0.12in; margin-bottom: 0.14in; }
.demo-box { background: #F8F8F8; padding: 0.14in; border-left: 0.05in solid #1E50D8; }
.demo-box-lbl  { font-size: 9pt; font-weight: 700; color: #1E50D8; text-transform: uppercase; margin-bottom: 4px; }
.demo-box-body { font-size: 10pt; color: #0A0A0A; line-height: 1.4; }
"""

# ── HTML slide definitions ────────────────────────────────────────────────────

slides_html = [
    (1, "Cover", """
<div class="cover-slide">
  <div class="cover-band">
    <div class="cover-logo"><span class="bizi">Bizi</span><span class="bid">Bid</span></div>
    <div class="cover-tagline">Barbados' First Retail Auction &amp; Deals Platform</div>
    <div class="cover-sub">Partnership Opportunity &nbsp;·&nbsp; Confidential</div>
  </div>
  <div class="cover-footer">
    Presented by INTERXDB &nbsp;·&nbsp; Digital Retail &amp; Auction Platform &nbsp;·&nbsp; Barbados<br>
    bizibid.com &nbsp;·&nbsp; contact@interxdb.com &nbsp;·&nbsp; Tel 1(246) 241-3771
  </div>
</div>"""),

    (2, "What is BiziBid?", """
<div class="slide-label">Platform Overview</div>
<div class="slide-h1">What is BiziBid?</div>
<div class="slide-sub">A mobile-first auction and deals platform built exclusively for Barbados' retail market.</div>
<div class="card-row">
  <div class="card"><div class="card-icon">★</div><div class="card-title">The BiziBid Mobile App</div><p>A free app (iOS &amp; Android) where registered customers browse live auctions and exclusive deals — 24 hours a day, 7 days a week.</p></div>
  <div class="card"><div class="card-icon">★</div><div class="card-title">Your Back-Office Portal</div><p>A private management portal at BiziBid.com. Upload products, create deals, review auction results, and track performance.</p></div>
  <div class="card"><div class="card-icon">★</div><div class="card-title">Fully Operated by INTERXDB</div><p>INTERXDB owns and operates every part of the platform — the app, the servers, the technology. You simply list your products. We handle everything else.</p></div>
</div>"""),

    (3, "The Buzz", """
<div class="cover-slide buzz-slide">
  <div style="padding:0.35in 0.5in 0;">
    <div class="buzz-logo"><span class="bizi">Bizi</span><span class="bid">Bid</span></div>
    <div style="margin-top:0.12in;"><span class="buzz-badge">SOMETHING NEW IS COMING TO BARBADOS</span></div>
    <div class="buzz-h1">Wake the Market.</div>
    <div class="buzz-sub">Your customers are ready for something exciting — are you first?</div>
    <div class="buzz-boxes">
      <div class="buzz-box"><div class="buzz-box-title">The buzz is building.</div><p>Barbados' retail customers are craving deals, excitement, and a reason to engage with their favourite stores again. BiziBid is that reason.</p></div>
      <div class="buzz-box"><div class="buzz-box-title">First mover wins.</div><p>Founding retail partners secure the most visible positions in the app before it opens to competitors. The stores that join first are the stores that dominate.</p></div>
      <div class="buzz-box"><div class="buzz-box-title">Your brand stays alive 24/7.</div><p>Even when your store doors are closed, your name, your products, and your deals are on every customer's phone — generating excitement around the clock.</p></div>
    </div>
    <div class="buzz-footer" style="margin-top:0.18in;">Early partners lock in rates and placement for 12 months. Positions are limited.</div>
  </div>
</div>"""),

    (4, "User Registration & Demographics", """
<div class="slide-label">Customer Intelligence</div>
<div class="slide-h1">Know Your Customers Before They Walk In</div>
<div class="slide-sub">Every BiziBid user registers with verified personal details — giving your store valuable demographic intelligence.</div>
<div class="with-phone">
  <div class="text-side">
    <ul class="big-bullets" style="margin-bottom:0.14in;">
      <li>Full name, verified email, and mobile number collected at registration</li>
      <li>Date of birth — every user must be 18 or older — you know your age demographic</li>
      <li>Parish / address — you see exactly which areas your customers come from</li>
      <li>All data is aggregated and visible in your back-office statistics dashboard</li>
      <li>Track which age groups bid most, which parishes are most active, and more</li>
      <li>No other advertising platform in Barbados gives retail stores this level of insight</li>
    </ul>
    <div class="demo-boxes">
      <div class="demo-box"><div class="demo-box-lbl">Age Captured</div><div class="demo-box-body">18+ verified. Know if your buyers are 18-30, 30-45, or 45+.</div></div>
      <div class="demo-box"><div class="demo-box-lbl">Parish Known</div><div class="demo-box-body">See which districts your most active bidders come from.</div></div>
      <div class="demo-box"><div class="demo-box-lbl">Real Contact</div><div class="demo-box-body">Verified mobile &amp; email — not anonymous traffic. Real people.</div></div>
    </div>
  </div>
  <div class="phone-wrap">
    <div class="phone-outer">
      <div class="phone-speaker"></div>
      <div class="phone-screen">
        <div class="ph-status"><span>9:41</span><span>||| 100%</span></div>
        <div class="ph-header">Create Account</div>
        <div class="ph-field"><div class="ph-field-lbl">First Name</div><div class="ph-field-val">John</div></div>
        <div class="ph-field"><div class="ph-field-lbl">Last Name</div><div class="ph-field-val">Husbands</div></div>
        <div class="ph-field"><div class="ph-field-lbl">Date of Birth (18+ required)</div><div class="ph-field-val">15 / 04 / 1992</div></div>
        <div class="ph-field"><div class="ph-field-lbl">Parish</div><div class="ph-field-val">St. Michael</div></div>
        <div class="ph-field"><div class="ph-field-lbl">Email Address</div><div class="ph-field-val">john@gmail.com</div></div>
        <div class="ph-field"><div class="ph-field-lbl">Mobile Number</div><div class="ph-field-val">+1 (246) 256-8821</div></div>
        <div class="ph-btn">REGISTER FREE</div>
      </div>
      <div class="phone-home-bar"></div>
    </div>
  </div>
</div>"""),

    (5, "Inside the App", """
<div class="slide-label">The BiziBid App</div>
<div class="slide-h1">Inside the App — Your Store on Every Screen</div>
<div class="slide-sub">The BiziBid home screen is a live keypad of retail stores — your logo always visible, always active.</div>
<div class="with-phone">
  <div class="text-side">
    <div class="card-row two" style="margin-bottom:0.12in;">
      <div class="card"><div class="card-title">Pinned Store</div><p>Customers pin their favourite store at the top — maximum visibility every time they open the app.</p></div>
      <div class="card"><div class="card-title">Live Badge</div><p>Colour-coded badge shows active bids and alerts on your store button — customers react immediately.</p></div>
      <div class="card"><div class="card-title">Store Keypad</div><p>Your logo is a direct-tap button. One press opens your full auction and deals screen instantly.</p></div>
      <div class="card"><div class="card-title">Deal Strip</div><p>Active Deals of the Week appear as a scrollable banner across the top of the feed.</p></div>
    </div>
  </div>
  <div class="phone-wrap">
    <div class="phone-outer">
      <div class="phone-speaker"></div>
      <div class="phone-screen">
        <div class="ph-status"><span>9:41</span><span>||| 100%</span></div>
        <div class="ph-header">BiziBid</div>
        <div class="deal-strip">DEAL OF THE WEEK: 32" Smart TV  $199</div>
        <div class="ph-label">YOUR STORES</div>
        <div class="keypad-grid">
          <div class="store-key"><div class="store-key-init">M</div><div class="store-key-name">Massy</div><span class="store-badge">5</span></div>
          <div class="store-key"><div class="store-key-init">C</div><div class="store-key-name">Carters</div><span class="store-badge green">2</span></div>
          <div class="store-key"><div class="store-key-init">H</div><div class="store-key-name">H&amp;B</div></div>
          <div class="store-key"><div class="store-key-init">P</div><div class="store-key-name">Price-R</div><span class="store-badge">1</span></div>
          <div class="store-key"><div class="store-key-init">C</div><div class="store-key-name">Courts</div></div>
          <div class="store-key"><div class="store-key-init">M</div><div class="store-key-name">MasonMrt</div><span class="store-badge green">3</span></div>
        </div>
        <div class="ph-label" style="margin-top:4px;">MY ACTIVE BIDS</div>
        <div class="bid-row top"><span class="bid-user">32" Smart TV</span><span class="bid-amt">$215</span></div>
        <div class="bid-row"><span class="bid-user">Gas BBQ Grill</span><span class="bid-amt">$88</span></div>
      </div>
      <div class="phone-home-bar"></div>
    </div>
  </div>
</div>"""),

    (6, "How Customers Bid", """
<div class="slide-label">The Bidding Experience</div>
<div class="slide-h1">How Customers Bid — Simple, Fast, Addictive</div>
<div class="slide-sub">Three taps and your customer is competing live. The simplicity is what drives repeat engagement.</div>
<div class="with-phone">
  <div class="text-side">
    <div class="step-row"><div class="step-num">1</div><div><div class="step-title">Open BiziBid</div><div class="step-body">Customer opens the app and sees your store on the home screen with an active bid badge.</div></div></div>
    <div class="step-row"><div class="step-num">2</div><div><div class="step-title">Tap Your Store</div><div class="step-body">One tap opens your auction listing — product image, description, current bid price, and live countdown.</div></div></div>
    <div class="step-row"><div class="step-num">3</div><div><div class="step-title">Tap BID NOW</div><div class="step-body">A single tap places their bid. They are instantly in the lead. No forms, no complicated process.</div></div></div>
    <div class="step-row"><div class="step-num">4</div><div><div class="step-title">Outbid Alert Fires</div><div class="step-body">If outbid, a push notification arrives immediately on their phone — bringing them straight back.</div></div></div>
    <div class="step-row"><div class="step-num">5</div><div><div class="step-title">Win &amp; Collect</div><div class="step-body">Auction closes. Winner receives a unique collection code. They bring it to your store to collect.</div></div></div>
  </div>
  <div class="phone-wrap">
    <div class="phone-outer">
      <div class="phone-speaker"></div>
      <div class="phone-screen">
        <div class="ph-status"><span>9:41</span><span>||| 100%</span></div>
        <div class="ph-header">Live Auction</div>
        <div class="ph-prod">
          <div class="ph-prod-name">32" Smart TV — Samsung</div>
          <div class="ph-prod-sub">Starting bid: $150.00</div>
          <div class="ph-prod-timer">Ends in  00:04:32</div>
        </div>
        <div class="ph-cur-bid">
          <div class="ph-cur-lbl">CURRENT BID</div>
          <div class="ph-cur-amt">$215.00</div>
        </div>
        <div class="ph-label">BID HISTORY</div>
        <div class="bid-row top"><span class="bid-user">j.husbands</span><span class="bid-amt">$215.00</span></div>
        <div class="bid-row"><span class="bid-user">r.alleyne</span><span class="bid-amt">$200.00</span></div>
        <div class="bid-row"><span class="bid-user">b.springer</span><span class="bid-amt">$185.00</span></div>
        <div class="ph-btn">BID NOW — $220.00</div>
      </div>
      <div class="phone-home-bar"></div>
    </div>
  </div>
</div>"""),

    (7, "Bid & Outbid Notifications", """
<div class="slide-label">Real-Time Notifications</div>
<div class="slide-h1">Notifications That Bring Customers Back Instantly</div>
<div class="slide-sub">Every bid event triggers a push notification — customers never miss a moment of the auction.</div>
<div class="with-phone">
  <div class="text-side">
    <div style="background:#F0FBF4;border-left:0.06in solid #16A34A;padding:0.16in;margin-bottom:0.14in;">
      <div style="font-size:10pt;font-weight:700;color:#16A34A;margin-bottom:6px;">WINNING BID NOTIFICATION</div>
      <p style="font-size:12pt;color:#4B5563;line-height:1.5;margin-bottom:8px;">Sent immediately when a customer takes the lead. The positive reinforcement is powerful — it keeps them emotionally invested and returning to the app to check their status.</p>
      <p style="font-size:11pt;font-weight:700;color:#0A0A0A;">"You're winning! Keep it up."</p>
    </div>
    <div style="background:#FDF2F2;border-left:0.06in solid #DC2626;padding:0.16in;">
      <div style="font-size:10pt;font-weight:700;color:#DC2626;margin-bottom:6px;">OUTBID NOTIFICATION</div>
      <p style="font-size:12pt;color:#4B5563;line-height:1.5;margin-bottom:8px;">Fires the instant another customer places a higher bid. Most counter-bids happen within 90 seconds. Customers tap straight back into the app to respond.</p>
      <p style="font-size:11pt;font-weight:700;color:#0A0A0A;">"You've been outbid! Tap to stay in the lead."</p>
    </div>
  </div>
  <div class="phone-wrap">
    <div class="phone-outer">
      <div class="phone-speaker"></div>
      <div class="phone-screen">
        <div class="ph-status"><span>9:41</span><span>||| 100%</span></div>
        <div class="ph-header">Notifications</div>
        <div class="ph-notif"><span class="ph-notif-time">now</span><div class="ph-notif-title">You're winning!</div><div class="ph-notif-body">32" Smart TV — bid $215</div></div>
        <div class="ph-notif"><span class="ph-notif-time">2 min</span><div class="ph-notif-title">Your bid is leading</div><div class="ph-notif-body">Gas BBQ Grill — bid $88</div></div>
        <div class="ph-notif red"><span class="ph-notif-time">5 min</span><div class="ph-notif-title">You've been outbid!</div><div class="ph-notif-body">32" Smart TV — someone bid $225</div></div>
        <div class="ph-notif red"><span class="ph-notif-time">5 min</span><div class="ph-notif-title">Tap to counter-bid</div><div class="ph-notif-body">You had $215 — bid $225 to lead</div></div>
        <div class="ph-notif amber"><span class="ph-notif-time">1 hr</span><div class="ph-notif-title">Auction ending soon!</div><div class="ph-notif-body">Gas BBQ Grill closes in 2 hours</div></div>
        <div class="ph-btn">VIEW AUCTIONS</div>
      </div>
      <div class="phone-home-bar"></div>
    </div>
  </div>
</div>"""),

    (8, "Deals of the Week & Month", """
<div class="slide-label">Feature — Deals</div>
<div class="slide-h1">Deals of the Week &amp; Month — Drive Foot Traffic</div>
<div class="slide-sub">Fixed-price deals showcase your best offers and bring customers through your doors.</div>
<div class="two-col">
  <ul class="big-bullets">
    <li>You set a fixed price — no bidding, no auction pressure</li>
    <li>Customers see your deal prominently in the app all week or all month</li>
    <li>Customers can Like the deal — like count is visible to all users (social proof)</li>
    <li>Customers share deals directly to WhatsApp, Instagram, and other platforms</li>
    <li>Shared links show the deal to people who do not yet have the app — driving new downloads</li>
    <li>Walk-in customers ask for the deal by name — increased foot traffic guaranteed</li>
    <li>Deals of the Month offer extended 30-day visibility at premium placement</li>
  </ul>
  <div style="display:flex;flex-direction:column;gap:0.12in;">
    <div style="background:#F8F8F8;padding:0.16in;border-left:0.05in solid #C9A84C;">
      <div style="font-size:12pt;font-weight:700;color:#0A0A0A;margin-bottom:6px;">Deal of the Week</div>
      <p style="font-size:11pt;color:#4B5563;line-height:1.6;">7-day fixed-price feature<br>Premium placement in app<br>Like + Share enabled<br>Drives walk-in purchase</p>
    </div>
    <div style="background:#0B1F3A;padding:0.16in;border-left:0.05in solid #C9A84C;">
      <div style="font-size:12pt;font-weight:700;color:#C9A84C;margin-bottom:6px;">Deal of the Month</div>
      <p style="font-size:11pt;color:#C0C8D8;line-height:1.6;">Full 30-day visibility<br>Top placement position<br>Maximum brand exposure<br>Ideal for key product lines</p>
    </div>
  </div>
</div>"""),

    (9, "A 24/7 Brand Portal", """
<div class="slide-label">Brand Strategy</div>
<div class="slide-h1">More Than a Marketplace — A 24/7 Brand Portal</div>
<div class="slide-sub">BiziBid is your store's permanent digital home. It keeps your brand alive in customers' minds every single day.</div>
<div class="card-row two" style="gap:0.12in;">
  <div class="card"><div class="card-title">Open 24/7</div><p>Your store is visible and active on BiziBid even when your doors are closed. Customers browse your products at midnight, on weekends, on public holidays.</p></div>
  <div class="card"><div class="card-title">Not just a sale — an experience</div><p>Customers follow your store, watch your auctions, share your deals, and talk about you. That is brand building — happening automatically.</p></div>
  <div class="card"><div class="card-title">A daily digital billboard</div><p>Your store logo sits on every registered customer's phone home screen. No billboard delivers that kind of daily reach at this price.</p></div>
  <div class="card"><div class="card-title">Customer loyalty, digitally</div><p>When a customer pins your store, sets a notification, and bids week after week — that is a loyal customer. BiziBid builds that relationship for you.</p></div>
  <div class="card"><div class="card-title">Your competitors are still offline</div><p>Most Barbados retail stores have no meaningful digital presence. BiziBid makes your store the digital leader in your category before your rivals even react.</p></div>
  <div class="card"><div class="card-title">The app is your voice</div><p>Every auction, every deal, every share is your brand speaking to customers. BiziBid gives your brand a voice that never goes silent.</p></div>
</div>"""),

    (10, "Live Auctions", """
<div class="slide-label">Feature — Auctions</div>
<div class="slide-h1">Live Auctions — Excitement That Sells</div>
<div class="slide-sub">Customers compete in real-time for your products. The excitement keeps them engaged and coming back.</div>
<div class="two-col">
  <ul class="big-bullets">
    <li>You upload products with a starting bid price (set below shelf price)</li>
    <li>Up to 25 items per week rotate through the platform automatically</li>
    <li>5 new items released each day at varied morning hours — customers check in daily</li>
    <li>Auctions run for 3, 7, or 14 days depending on the item</li>
    <li>Real-time bidding — customers see bids update live on their screen</li>
    <li>Unsold items re-enter the rotation for the following week automatically</li>
  </ul>
  <div class="dark-box">
    <div class="dark-box-title">When a Customer Wins</div>
    <ol class="win-steps">
      <li>Auction closes — winner is confirmed</li>
      <li>A unique collection code is generated</li>
      <li>Winner receives code on their phone</li>
      <li>Winner visits your store and shows code to cashier</li>
      <li>Cashier confirms the code — item is collected</li>
      <li>Record of the sale is stored in your portal</li>
    </ol>
    <div class="dark-note">No online payment needed — all purchases are cash in-store.</div>
  </div>
</div>"""),

    (11, "Back-Office Portal", """
<div class="slide-label">Your Back-Office — BiziBid.com</div>
<div class="slide-h1">Your Private Management Portal</div>
<div class="slide-sub">Log in at BiziBid.com to manage your entire BiziBid presence — available 24/7 from any browser.</div>
<div class="card-row" style="gap:0.12in;">
  <div class="card sm"><div class="card-title">Bulk Product Upload</div><p>Upload multiple products at once with images, descriptions, and starting bid price. No technical knowledge required.</p></div>
  <div class="card sm"><div class="card-title">Deals Management</div><p>Create and schedule Deals of the Week and Deals of the Month. Set your own price and valid dates. Activate with one click.</p></div>
  <div class="card sm"><div class="card-title">Auction Results</div><p>View all completed auctions — winning bid, customer name, collection code. Full record of every sale.</p></div>
  <div class="card sm"><div class="card-title">Collection Codes</div><p>See the unique code for each auction winner so your cashier can verify the customer quickly and confidently.</p></div>
  <div class="card sm"><div class="card-title">Invoice Download</div><p>View and download your monthly platform invoices as PDF. Full record of your subscription and activity charges.</p></div>
  <div class="card sm"><div class="card-title">Store Statistics</div><p>Live dashboard — user engagement, bid activity, deal shares, and more. Know exactly how customers interact with your store.</p></div>
</div>"""),

    (12, "Store Statistics", """
<div class="slide-label">Back-Office — Statistics Dashboard</div>
<div class="slide-h1">Know Your Customers — Live Store Statistics</div>
<div class="slide-sub">Your back-office portal shows you exactly how registered BiziBid users interact with your store in real time.</div>
<div class="stat-row">
  <div class="stat-card"><div class="stat-num">1,240</div><div class="stat-lbl">Registered users following your store</div></div>
  <div class="stat-card"><div class="stat-num">86</div><div class="stat-lbl">Active bids on your items today</div></div>
  <div class="stat-card"><div class="stat-num">312</div><div class="stat-lbl">Users watching your products</div></div>
  <div class="stat-card"><div class="stat-num">94%</div><div class="stat-lbl">Auction completion rate this month</div></div>
</div>
<div class="stat-row" style="margin-top:0.1in;">
  <div class="stat-card"><div class="stat-num">48</div><div class="stat-lbl">Deals shared on WhatsApp &amp; Instagram</div></div>
  <div class="stat-card"><div class="stat-num">$8,400</div><div class="stat-lbl">Total auction value generated this month</div></div>
  <div class="stat-card"><div class="stat-num">22</div><div class="stat-lbl">New app users from your shared links</div></div>
  <div class="stat-card"><div class="stat-num">4.8★</div><div class="stat-lbl">Average customer engagement score</div></div>
</div>
<div class="footnote-box" style="margin-top:0.12in;">All statistics are updated in real time in your back-office portal. See age group breakdowns and which parishes your most active bidders come from. No data is shared with other retail partners.</div>"""),

    (13, "10x to 100x Visibility", """
<div class="slide-label">Viral Sharing Power</div>
<div class="slide-h1">One Share = 10x to 100x Your Visibility</div>
<div class="slide-sub">When a customer shares your product on WhatsApp or social media, your brand reaches their entire network instantly.</div>
<div class="share-boxes">
  <div class="share-box"><div class="share-mult">10x</div><div class="share-ch">WhatsApp Message</div><div class="share-body">Sent to a group of 50-150 people. Everyone in that group sees your store name, product, and price.</div></div>
  <div class="share-box"><div class="share-mult">50x</div><div class="share-ch">WhatsApp Status</div><div class="share-body">Visible to all of a customer's WhatsApp contacts for 24 hours. Average reach: 200-400 people per post.</div></div>
  <div class="share-box"><div class="share-mult">100x</div><div class="share-ch">Facebook / Instagram</div><div class="share-body">Liked, shared, or reposted content reaches friends-of-friends. A single viral share can reach thousands.</div></div>
</div>
<div class="share-footer">
  <div class="share-footer-h">People love to share a great deal.</div>
  <p>A customer who finds a $50 item auction starting at $12 will share it to every WhatsApp group they have. That is free marketing to hundreds of people who have never heard of your store — and it costs you nothing. Every shared link drives new app downloads, and new app users are new bidders on your products.</p>
</div>"""),

    (14, "Viral Growth Funnel", """
<div class="slide-label">Feature — Viral Growth</div>
<div class="slide-h1">Every Share is Free Marketing for Your Store</div>
<div class="slide-sub">When customers share your auctions and deals, your store name reaches people who have never heard of BiziBid.</div>
<div class="flow-row">
  <div class="flow-box dark"><span class="flow-num">01</span>Your customer sees a great deal or live auction</div>
  <div class="flow-arrow">-&gt;</div>
  <div class="flow-box light"><span class="flow-num gold">02</span>They share it to WhatsApp or Instagram</div>
  <div class="flow-arrow">-&gt;</div>
  <div class="flow-box dark"><span class="flow-num">03</span>Their contact sees your store name, price &amp; product</div>
  <div class="flow-arrow">-&gt;</div>
  <div class="flow-box light"><span class="flow-num gold">04</span>They visit BiziBid.com and download the app</div>
  <div class="flow-arrow">-&gt;</div>
  <div class="flow-box dark"><span class="flow-num">05</span>New registered bidder — on your auctions</div>
</div>
<div class="footnote-box">
  <b>Auction shares</b> show the live bid price and countdown — creating instant urgency for the viewer.<br>
  <b>Deal shares</b> show your store name, product image, and fixed price — driving walk-in intent.<br>
  All shared links direct visitors to BiziBid.com to download the app. New downloads = new bidders on your products.
</div>"""),

    (15, "Keypad Placement", """
<div class="slide-label">Placement &amp; Visibility</div>
<div class="slide-h1">Secure Your Position on the BiziBid Keypad</div>
<div class="slide-sub">Your store logo placement on the app home screen is a paid position. First come, first placed.</div>
<div class="tier-row">
  <div class="tier gold-tier"><div class="tier-label">Tier 1 — Pinned Position</div><p>Top of every user's home screen. Your logo is displayed largest, with full notification badge visibility. One store only — the most premium placement available.</p></div>
  <div class="tier navy-tier"><div class="tier-label">Tier 2 — Primary Keypad</div><p>First 6–9 visible positions in the main keypad grid. Seen immediately when the app opens. No scrolling required. Highest traffic positions after Tier 1.</p></div>
  <div class="tier grey-tier"><div class="tier-label">Tier 3 — Slide-Out Panel</div><p>Your store appears in the extended keypad panel. Still fully accessible and notified — ideal for entry-level platform participation.</p></div>
  <div class="tier dark-tier"><div class="tier-label">Banner Advertisements</div><p>A rotating advertisement strip displayed across all screens. Separate from keypad placement — available to any registered partner.</p></div>
</div>"""),

    (16, "Pricing", """
<div class="slide-label">Partnership Pricing</div>
<div class="slide-h1">Simple, Transparent Pricing (BBD)</div>
<div class="slide-sub">All rates are fixed for the first 12 months for founding partners. No hidden fees.</div>
<div class="offer-box"><span class="offer-days">60 DAYS FREE</span> No monthly cost for the first 60 days. A <b>$500.00 BBD</b> deposit secures your primary advertisement placement — guaranteeing your store is featured prominently to all registered users from day one.</div>
<table class="price-table">
  <thead><tr><th>Service</th><th>Description</th><th>Rate (BBD)</th></tr></thead>
  <tbody>
    <tr><td><b>Base Listing Package</b></td><td>Up to 25 items listed per week on the app</td><td class="price">$450.00 / month</td></tr>
    <tr class="alt"><td><b>Additional Items</b></td><td>Per additional 10 items above the 25-item base</td><td class="price">$150.00 / month</td></tr>
    <tr><td><b>Deal of the Week</b></td><td>Featured deal highlighted to all users for 7 days</td><td class="price">$150.00 / week</td></tr>
    <tr class="alt"><td><b>Deal of the Month</b></td><td>Premium featured deal — top placement for full month</td><td class="price">$200.00 / month</td></tr>
  </tbody>
</table>
<div class="table-note">All pricing in Barbados Dollars (BBD). Rates fixed for 12 months for founding partners. Keypad placement and banner ad pricing available on request.</div>"""),

    (17, "Why BiziBid?", """
<div class="slide-label">Why BiziBid?</div>
<div class="slide-h1">Built for Barbados Retail. Starting Today.</div>
<div class="card-row" style="margin-top:0.15in;">
  <div class="card"><div class="card-title">Barbados First</div><p>Built specifically for the Barbados market — the culture, the community, the retail landscape. Not a foreign platform adapted for the island.</p></div>
  <div class="card"><div class="card-title">Zero Technical Effort</div><p>You upload products and set prices. INTERXDB operates the entire platform. Nothing more is required from you.</p></div>
  <div class="card"><div class="card-title">Keep Your Store Relevant</div><p>BiziBid gives your store a daily digital presence that keeps you competitive — without the cost of building your own platform.</p></div>
  <div class="card"><div class="card-title">Exclusive &amp; Professional</div><p>Only verified, registered businesses may list products. No individual sellers. A premium marketplace that reflects well on your brand.</p></div>
  <div class="card"><div class="card-title">Your Target Market</div><p>BiziBid targets medium-to-high income earners and family units — the customers who already shop at your store.</p></div>
  <div class="card"><div class="card-title">Founding Partner Advantage</div><p>Founding partners lock in rates for 12 months and secure preferred keypad placement before the platform opens to all retailers.</p></div>
</div>"""),

    (18, "Get Started", """
<div class="cover-slide cta-slide">
  <div class="cover-band">
    <div class="cover-logo" style="font-size:42px;"><span class="bizi">Bizi</span><span class="bid">Bid</span></div>
    <div class="cover-tagline" style="font-size:26px;margin-top:10px;">Ready to Join BiziBid?</div>
    <div class="cover-sub">Secure your founding partner status today — positions are limited.</div>
    <div class="cta-cards">
      <div class="cta-card"><div class="cta-icon">P</div><div class="cta-title">Schedule a Presentation</div><p>We come to you. A personalised walkthrough at your convenience.</p></div>
      <div class="cta-card"><div class="cta-icon">W</div><div class="cta-title">View Your Interactive Demo</div><p>A personalised before-and-after demonstration prepared for your store.</p></div>
      <div class="cta-card"><div class="cta-icon">E</div><div class="cta-title">Contact Us Directly</div><p>Our team is ready to answer any questions and walk you through the terms.</p></div>
    </div>
  </div>
  <div class="cover-footer">
    BiziBid &nbsp;·&nbsp; Operated by INTERXDB &nbsp;·&nbsp; contact@interxdb.com &nbsp;·&nbsp; Tel 1(246) 241-3771 &nbsp;·&nbsp; BiziBid.com<br>
    <span style="font-size:10px;color:#6B7280;">Confidential — For Addressee Only</span>
  </div>
</div>"""),
]

# ── Assemble HTML ─────────────────────────────────────────────────────────────
html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>BiziBid Partner Presentation</title>
<style>{css}</style>
</head>
<body>
"""

for num, title, content in slides_html:
    is_cover = title in ("Cover", "Get Started")
    if is_cover:
        html += f'<div class="slide" style="padding:0;">\n  {content}\n  <div class="page-num">{num} / {TOTAL}</div>\n</div>\n'
    else:
        html += f"""<div class="slide">
  <div class="logo-top"><span class="bizi">Bizi</span><span class="bid">Bid</span></div>
  {content}
  <div class="gold-bar"></div>
  <div class="page-num">{num} / {TOTAL}</div>
</div>
"""

html += "</body></html>"

pdf_path = os.path.join(out_dir, "BiziBid-Partner-Presentation-PRINT.html")
with open(pdf_path, "w", encoding="utf-8") as f:
    f.write(html)
print(f"Saved: {pdf_path}")
print("Open in Chrome > Print > Save as PDF (paper size 13.33 x 7.5in, landscape)")
print("Done. 18 slides generated.")

